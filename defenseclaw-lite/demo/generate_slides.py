#!/usr/bin/env python3
"""Generate DefenseClaw Lite demo presentation — 3 slides."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
CISCO_BLUE = RGBColor(0x04, 0x9F, 0xD9)
RED_BLOCK = RGBColor(0xE8, 0x3E, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE = RGBColor(0xFF, 0xAA, 0x00)
CISCO_GREEN = RGBColor(0x6C, 0xC0, 0x4A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_dark_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BG
    return slide


def add_text(slide, text, top, left=Inches(0.8), width=Inches(11.5),
             size=Pt(44), color=WHITE, bold=True):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(6)
    return txBox


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = add_dark_slide()
add_text(slide, "DefenseClaw Lite", Inches(1.8), size=Pt(60), color=CISCO_BLUE)
add_text(slide, "AI Agent Security for IoT & Edge Devices", Inches(3.2), size=Pt(32), color=WHITE, bold=False)
add_text(slide,
         "54 KB  ·  <3μs decisions  ·  25 KB RAM  ·  Open Source\n\n"
         "Protects AI agents on Raspberry Pi, Jetson, industrial gateways,\n"
         "robots, and any Linux device.",
         Inches(4.5), size=Pt(20), color=LIGHT_GRAY, bold=False)
add_text(slide, "Cisco AI Defense", Inches(6.5), size=Pt(18), color=CISCO_BLUE, bold=False)

# ============================================================
# SLIDE 2: Architecture
# ============================================================
slide = add_dark_slide()
add_text(slide, "How It Works", Inches(0.4), size=Pt(44), color=CISCO_BLUE)

add_text(slide,
         "┌─────────────────────────────────────────────────────────────────┐\n"
         "│  AI Agent (PicoClaw)                                            │\n"
         "│       │                                                         │\n"
         "│       │ every tool call                                         │\n"
         "│       ▼                                                         │\n"
         "│  ┌──────────────────────────────────────────────────────────┐  │\n"
         "│  │  DefenseClaw Lite (54KB, <3μs)                           │  │\n"
         "│  │                                                          │  │\n"
         "│  │  1. Input validation     5. Sequence detection           │  │\n"
         "│  │  2. Rate limiting        6. Verdict cache                │  │\n"
         "│  │  3. Hash deny-list       7. Cloud escalation             │  │\n"
         "│  │  4. Dest filtering                                       │  │\n"
         "│  │                                                          │  │\n"
         "│  │  → ALLOW (sensors, reads)    → BLOCK (exec, actuate)    │  │\n"
         "│  └──────────────────────────────────────────────────────────┘  │\n"
         "│       │                                                         │\n"
         "│       ▼  only if ALLOWED                                       │\n"
         "│  Tool executes (robot moves, sensor reads, API calls)          │\n"
         "└─────────────────────────────────────────────────────────────────┘",
         Inches(1.4), size=Pt(13), color=LIGHT_GRAY, bold=False, left=Inches(1.5), width=Inches(10.5))

# ============================================================
# SLIDE 3: Closing
# ============================================================
slide = add_dark_slide()
add_text(slide, "DefenseClaw Lite", Inches(2.0), size=Pt(54), color=CISCO_BLUE)
add_text(slide,
         "The IoT AI agent era is here. Now, so is the security.",
         Inches(3.5), size=Pt(28), color=WHITE, bold=False)
add_text(slide,
         "github.com/cisco-ai-defense/defenseclaw\n\n"
         "76 KB binary  ·  2.7μs latency  ·  370K decisions/sec  ·  Zero dependencies\n"
         "Supports: PicoClaw  |  Roadmap: Bubbaloop, IoT-Edge-MCP, TinyAgent",
         Inches(4.8), size=Pt(18), color=LIGHT_GRAY, bold=False)
add_text(slide, "Cisco AI Defense", Inches(6.5), size=Pt(18), color=CISCO_BLUE, bold=False)

# ============================================================
output_path = "/Users/nghodki/workspace/defenseclaw-workspace/defenseclaw/defenseclaw-lite/demo/DefenseClaw-Lite-Demo.pptx"
prs.save(output_path)
print(f"Saved: {output_path} ({len(prs.slides)} slides)")
